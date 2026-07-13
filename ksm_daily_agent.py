"""
KSM Marketing Agent — ежедневная автоматизация
Пайплайн: Roistat -> Claude (Anthropic API) -> хранилище -> Telegram

Скрипт запускается по расписанию (GitHub Actions, cron '0 1 * * *' = 07:00
по Алматы) и один раз за вызов делает:
  - ежедневное сравнение (вчера / позавчера) — всегда;
  - еженедельное сравнение (среда → вторник, против такой же недели ранее) —
    только по средам, когда прошедшая неделя только что полностью закрылась.

Секреты нигде не хардкожены — только через переменные окружения (см.
.env.example рядом). Перед запуском: pip install requests python-dotenv
python-pptx python-docx openpyxl
"""

import os
import sys
import json
import base64
import datetime
from datetime import timezone, timedelta

import requests
from dotenv import load_dotenv

load_dotenv()

# ---------- Конфигурация из переменных окружения ----------
ROISTAT_API_KEY = os.environ["ROISTAT_API_KEY"]
ROISTAT_PROJECT_ID = os.environ["ROISTAT_PROJECT_ID"]
ANTHROPIC_API_KEY = os.environ["ANTHROPIC_API_KEY"]
TELEGRAM_BOT_TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN")
# Несколько получателей — через запятую в одном секрете (например
# "5313414485,987654321"), каждый должен один раз написать боту лично
# (иначе Telegram не разрешит боту писать первым — правило платформы).
# Список из TELEGRAM_CHAT_ID — это "ручной" базовый список; sync_telegram_recipients()
# ниже дополняет его теми, кто сам написал боту и назвал кодовое слово.
TELEGRAM_CHAT_IDS = [c.strip() for c in os.environ.get("TELEGRAM_CHAT_ID", "").split(",") if c.strip()]
# Кодовое слово для самостоятельной подписки в боте — секрет (репозиторий
# публичный, поэтому слово нельзя хардкодить в коде/config.json, иначе любой
# читающий репозиторий узнает его и сможет подписаться сам).
TELEGRAM_JOIN_CODE = os.environ.get("TELEGRAM_JOIN_CODE", "").strip().lower()

# Для ежедневных/еженедельных отчётов достаточно Sonnet.
CLAUDE_MODEL = "claude-sonnet-5"
CLAUDE_MAX_TOKENS = 8192

DAILY_SNAPSHOT_PATH = "latest_snapshot.json"
WEEKLY_SNAPSHOT_PATH = "latest_weekly_snapshot.json"
CONFIG_PATH = "config.json"
KNOWLEDGE_BASE_PATH = "knowledge_base.md"
KNOWLEDGE_DIR = "knowledge"

# Пороги и план по умолчанию — используются, если нет config.json и не заданы
# соответствующие переменные окружения. Значения совпадают с тем, что раньше
# было захардкожено в SETTINGS на сайте, чтобы поведение не поменялось молча.
DEFAULT_CONFIG = {
    "cplAlertPct": 15,
    "leadsDropAlertPct": 15,
    "visitsDropAlertPct": 25,
    "ctrDropAlertPct": 25,
    "cpcRiseAlertPct": 25,
    "conversionDropAlertPct": 25,
    "noLeadsSpendAlertTenge": 3000,
    "monthlyLeadsTarget": 400,
    "monthlyBudgetTarget": 8000000,
}


def load_config():
    """config.json — общий файл настроек, который также редактирует вкладка
    «Настройки» на сайте (через GitHub Contents API), чтобы пороги и план не
    расходились между сайтом и тем, что реально проверяет скрипт."""
    cfg = dict(DEFAULT_CONFIG)
    if os.path.exists(CONFIG_PATH):
        try:
            with open(CONFIG_PATH, encoding="utf-8") as f:
                cfg.update(json.load(f))
        except (json.JSONDecodeError, OSError) as e:
            print(f"[config.json] не удалось прочитать, использую значения по умолчанию: {e}", file=sys.stderr)
    # Переменные окружения, если заданы явно, имеют приоритет (обратная совместимость).
    if os.environ.get("CPL_ALERT_PCT"):
        cfg["cplAlertPct"] = float(os.environ["CPL_ALERT_PCT"])
    if os.environ.get("LEADS_DROP_ALERT_PCT"):
        cfg["leadsDropAlertPct"] = float(os.environ["LEADS_DROP_ALERT_PCT"])
    return cfg


CONFIG = load_config()

# Часовой пояс проекта в Roistat / компании — Алматы, UTC+6, без перехода на
# летнее время. ВАЖНО: раньше даты считались через datetime.date.today(),
# которая на раннере GitHub Actions берёт UTC, а не алматинское время. При
# ручном запуске поздно вечером (18:00-23:59 UTC = 00:00-05:59 по Алматы)
# это давало даты на 1 день раньше нужных — подтверждено на практике.
ALMATY_TZ = timezone(timedelta(hours=6))
ROISTAT_TZ_OFFSET = "+06:00"


def almaty_today():
    return datetime.datetime.now(ALMATY_TZ).date()


# ---------- 1. Roistat: получение данных ----------
# Подтверждено эмпирически: dimensions=["marker_level_1","marker_level_2"]
# даёт и канал верхнего уровня, и под-источник (кампания/тип размещения —
# то, как это делит сам Roistat, зависит от канала). metrics расширены
# показами и кликами — Roistat отдаёт их по факту (проверено), это позволяет
# считать CTR/CPC без отдельной интеграции с Google Ads/Яндекс.Директ.
def fetch_roistat_data(date_from, date_to):
    url = "https://cloud.roistat.com/api/v1/project/analytics/data"
    headers = {"Api-key": ROISTAT_API_KEY}
    params = {"project": ROISTAT_PROJECT_ID}
    body = {
        "dimensions": ["marker_level_1", "marker_level_2"],
        "metrics": ["visits", "leads", "marketing_cost", "impressions", "clicks"],
        "period": {
            "from": f"{date_from}T00:00:00{ROISTAT_TZ_OFFSET}",
            "to": f"{date_to}T23:59:59{ROISTAT_TZ_OFFSET}",
        },
    }
    resp = requests.post(url, headers=headers, params=params, json=body, timeout=30)
    if resp.status_code >= 400:
        print(f"Roistat вернул ошибку, статус {resp.status_code}.", file=sys.stderr)
        print("Тело ответа:", resp.text[:1000], file=sys.stderr)
        resp.raise_for_status()
    try:
        payload = resp.json()
    except requests.exceptions.JSONDecodeError:
        print("=" * 60, file=sys.stderr)
        print(f"Roistat ответил статусом {resp.status_code}, но это не JSON.", file=sys.stderr)
        print("Заголовок ответа (Content-Type):", resp.headers.get("Content-Type"), file=sys.stderr)
        print("Первые 500 символов ответа:", repr(resp.text[:500]), file=sys.stderr)
        print("=" * 60, file=sys.stderr)
        raise
    if payload.get("status") == "error":
        # Roistat заворачивает ошибки авторизации/параметров в HTTP 200 —
        # без этой проверки extract_totals/extract_channels молча вернут нули.
        raise RuntimeError(
            f"Roistat API вернул ошибку: {payload.get('description') or payload.get('error')}"
        )
    return payload


def _iter_items(roistat_response):
    data = roistat_response.get("data")
    if isinstance(data, list) and data:
        block = data[0]
    elif isinstance(data, dict):
        block = data
    else:
        block = {}
    return block.get("items") or []


def _metric_dict(item):
    m = {}
    for metric in item.get("metrics", []):
        name = metric.get("metric_name")
        m[name] = (m.get(name) or 0) + (metric.get("value") or 0)
    return m


def extract_totals(roistat_response):
    """Суммирует показатели по ВСЕМ каналам за период."""
    items = _iter_items(roistat_response)
    leads = cost = visits = impressions = clicks = 0
    for item in items:
        m = _metric_dict(item)
        leads += m.get("leads", 0)
        cost += m.get("marketing_cost", 0)
        visits += m.get("visits", 0)
        impressions += m.get("impressions", 0)
        clicks += m.get("clicks", 0)

    totals = {
        "leads": leads,
        "cost": cost,
        "visits": visits,
        "impressions": impressions,
        "clicks": clicks,
        "cpl": (cost / leads) if leads else None,
        # ВАЖНО: подтверждено на 3 независимых сверках с реальным дашбордом
        # Roistat (точное совпадение до сотых процента/до тенге) — сам Roistat
        # считает CTR и CPC через "visits", а не через "clicks". "clicks" —
        # это сырой клик по объявлению у рекламной системы, "visits" — это
        # уже отфильтрованный Roistat трекнутый заход на сайт; в дашборде
        # используется именно второе. Если считать через clicks (как было
        # раньше), цифры расходятся с дашбордом на 20-30%.
        "ctr": (visits / impressions * 100) if impressions else None,
        "cpc": (cost / visits) if visits else None,
        "cr": (leads / visits * 100) if visits else None,
    }
    # На stderr, а не stdout: инструменты вроде roistat_debug.py печатают
    # машиночитаемый JSON в stdout, диагностика не должна в него подмешиваться.
    print(f"[Roistat] строк обработано: {len(items)}, лиды={leads}, расход={cost:.2f}, визиты={visits}", file=sys.stderr)
    if items and leads == 0 and cost == 0 and visits == 0:
        print("Суммы всё равно нулевые — вот пример одной строки целиком, чтобы найти точную причину:", file=sys.stderr)
        print(json.dumps(items[0], ensure_ascii=False), file=sys.stderr)
    return totals


def _normalize_key_part(s):
    return (s or "").strip().lower()


def channel_key(marker1_value, marker1_title):
    """Стабильный ключ канала для накопительной истории — по «сырому»
    значению Roistat (marker_level_1.value), а не по отображаемому title:
    названия могут отличаться регистром/пробелами между запусками, а
    значение — нет. config.json может содержать "channelAliases" для
    ручного схлопывания настоящих дублей (например, если один и тот же
    канал у Roistat иногда приходит под двумя разными value)."""
    raw = marker1_value or marker1_title or ""
    key = _normalize_key_part(raw)
    aliases = CONFIG.get("channelAliases", {})
    return aliases.get(key, key)


def extract_channels(roistat_response):
    """Детерминированная (не зависящая от вольного текста Claude) разбивка
    по каналам верхнего уровня, с под-источниками (marker_level_2)."""
    items = _iter_items(roistat_response)
    channels = {}
    for item in items:
        dims = item.get("dimensions") or {}
        l1_dim = dims.get("marker_level_1") or {}
        l1 = l1_dim.get("title") or "Без источника"
        l1_key = channel_key(l1_dim.get("value"), l1)
        l2 = (dims.get("marker_level_2") or {}).get("title")
        m = _metric_dict(item)

        ch = channels.setdefault(l1_key, {
            "key": l1_key,
            "channel": l1,
            "visits": 0,
            "leads": 0,
            "cost": 0,
            "impressions": 0,
            "clicks": 0,
            "sub": [],
        })
        ch["visits"] += m.get("visits", 0)
        ch["leads"] += m.get("leads", 0)
        ch["cost"] += m.get("marketing_cost", 0)
        ch["impressions"] += m.get("impressions", 0)
        ch["clicks"] += m.get("clicks", 0)
        if l2:
            ch["sub"].append({
                "name": l2,
                "visits": m.get("visits", 0),
                "leads": m.get("leads", 0),
                "cost": m.get("marketing_cost", 0),
            })

    result = []
    for ch in channels.values():
        ch["cpl"] = (ch["cost"] / ch["leads"]) if ch["leads"] else None
        # См. комментарий в extract_totals() — CTR/CPC считаются через visits,
        # как в самом дашборде Roistat, а не через clicks.
        ch["ctr"] = (ch["visits"] / ch["impressions"] * 100) if ch["impressions"] else None
        ch["cpc"] = (ch["cost"] / ch["visits"]) if ch["visits"] else None
        result.append(ch)
    result.sort(key=lambda c: c["cost"], reverse=True)
    return result


def rank_channels(channels):
    """Лучший/худший канал по CPL (среди тех, где были лиды) + список
    каналов, которые тратят бюджет без единого лида."""
    with_leads = [c for c in channels if c["leads"] > 0 and c["cpl"] is not None]
    best = min(with_leads, key=lambda c: c["cpl"]) if with_leads else None
    worst = max(with_leads, key=lambda c: c["cpl"]) if with_leads else None
    no_leads_spend = [
        c for c in channels
        if c["leads"] == 0 and c["cost"] >= CONFIG["noLeadsSpendAlertTenge"]
    ]
    return best, worst, no_leads_spend


# ---------- 1b. Накопительная история по каналам ----------
# Раньше анализ строился ТОЛЬКО на сравнении двух соседних дней/недель — не
# было накопленной истории, поэтому Claude сам честно отмечал, что выводы по
# 1-2 лидам статистически ненадёжны. Здесь — простое (без fine-tuning, это
# обычный вызов API без состояния) статистическое "обучение": скользящее
# окно сырых дневных цифр на канал в data/channel_history.json, из которого
# каждый запуск заново считаются базовые линии (среднее/std/тренд/z-score).
# Формулы и лимиты роста файла проверены отдельным дизайн-разбором с двумя
# независимыми судьями (fixed 90-дневное окно, без EWMA — см. обоснование
# в PR/чате: EWMA даёт неожиданно большой вес недавним точкам и может
# разойтись с исходными числами, а окно — нет, пересчитывается с нуля).
DATA_DIR = "data"
HISTORY_PATH = os.path.join(DATA_DIR, "channel_history.json")
HISTORY_WINDOW_DAYS = 90
INACTIVE_PRUNE_DAYS = 90
HARD_CHANNEL_CAP = 60
# Roistat пересматривает уже отданные цифры задним числом (подтверждено на
# практике: лиды за один и тот же день менялись между двумя проверками в
# пределах одной сессии, при неизменных визитах/расходе). Однократно забрать
# день и больше никогда не перепроверять — значит навсегда заморозить в
# истории устаревшую цифру. Поэтому каждый ежедневный запуск перезабирает
# не только "вчера", а последние REFRESH_WINDOW_DAYS дней целиком.
REFRESH_WINDOW_DAYS = 7
BASELINE_MIN_N = 14
STD_RELATIVE_FLOOR = 0.20
STD_ABS_EPSILON = 1e-6
HISTORY_METRICS = ("visits", "leads", "cost", "impressions", "clicks")


def _empty_history():
    return {
        "schema_version": 1,
        "window_size_days": HISTORY_WINDOW_DAYS,
        "inactive_prune_days": INACTIVE_PRUNE_DAYS,
        "hard_channel_cap": HARD_CHANNEL_CAP,
        "channels": {},
        "last_updated": None,
    }


def load_history():
    if not os.path.exists(HISTORY_PATH):
        return _empty_history()
    try:
        with open(HISTORY_PATH, encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError) as e:
        print(f"[История] Не удалось прочитать {HISTORY_PATH}, начинаю с чистого листа: {e}")
        return _empty_history()


def save_history(history):
    os.makedirs(DATA_DIR, exist_ok=True)
    history["last_updated"] = datetime.datetime.now(timezone.utc).isoformat()
    # Без indent — компактная сериализация вдвое меньше форматированной;
    # для ручной отладки: python -m json.tool data/channel_history.json
    with open(HISTORY_PATH, "w", encoding="utf-8") as f:
        json.dump(history, f, separators=(",", ":"), ensure_ascii=False, sort_keys=True)
    size = os.path.getsize(HISTORY_PATH)
    print(f"[История] Сохранено, размер файла: {size / 1024:.1f} КБ")
    if size > 2_000_000:
        # Предохранитель на случай, если hard_channel_cap не сработал как ожидалось.
        raise RuntimeError(
            f"data/channel_history.json превысил 2 МБ ({size} байт) — нужно расследовать рост числа каналов."
        )


def append_daily_record(history, date_str, channels):
    """Добавляет запись за date_str по каждому каналу; идемпотентно (повторный
    запуск за тот же день перезаписывает, а не дублирует), с обрезкой окна и
    прунингом каналов, переставших получать трафик."""
    for ch in channels:
        entry = history["channels"].setdefault(ch["key"], {"history": []})
        hist = entry["history"]
        hist[:] = [r for r in hist if r["date"] != date_str]
        hist.append({
            "date": date_str,
            "visits": ch["visits"],
            "leads": ch["leads"],
            "cost": ch["cost"],
            "impressions": ch["impressions"],
            "clicks": ch["clicks"],
        })
        hist.sort(key=lambda r: r["date"])
        if len(hist) > HISTORY_WINDOW_DAYS:
            del hist[0:len(hist) - HISTORY_WINDOW_DAYS]

    cutoff = (datetime.date.fromisoformat(date_str) - timedelta(days=INACTIVE_PRUNE_DAYS)).isoformat()
    for key in list(history["channels"].keys()):
        hist = history["channels"][key]["history"]
        if not hist or hist[-1]["date"] < cutoff:
            del history["channels"][key]

    if len(history["channels"]) > HARD_CHANNEL_CAP:
        def cost_28d(key):
            return sum(r["cost"] for r in history["channels"][key]["history"][-28:])
        ranked = sorted(history["channels"].keys(), key=cost_28d, reverse=True)
        dropped = ranked[HARD_CHANNEL_CAP:]
        for key in dropped:
            del history["channels"][key]
        if dropped:
            print(f"[История] Превышен лимит {HARD_CHANNEL_CAP} каналов, отброшено {len(dropped)}: {dropped}")

    return history


def refresh_recent_history(history, today, days_back=REFRESH_WINDOW_DAYS):
    """Перезабирает последние days_back дней из Roistat и перезаписывает их
    в истории (append_daily_record идемпотентен по дате — старая запись
    удаляется, новая добавляется). Это и есть механизм, которым история
    "узнаёт" о задним числом пересмотренных Roistat лидах/цифрах — без
    этого шага один раз сохранённый день навсегда остался бы устаревшим,
    даже если сам Roistat его потом поменяет."""
    for i in range(1, days_back + 1):
        d = today - timedelta(days=i)
        try:
            raw = fetch_roistat_data(str(d), str(d))
            channels = extract_channels(raw)
            append_daily_record(history, str(d), channels)
        except Exception as e:
            print(f"[История] Не удалось перепроверить {d}: {e}", file=sys.stderr)
    return history


def _mean(values):
    return sum(values) / len(values) if values else None


def _sample_std(values):
    n = len(values)
    if n < 2:
        return None
    m = _mean(values)
    return (sum((v - m) ** 2 for v in values) / (n - 1)) ** 0.5


def compute_baselines(history, exclude_dates):
    """Базовые линии по каждому каналу из уже накопленной истории, СТРОГО
    исключая exclude_dates (обычно сам анализируемый период) — иначе
    сравнение размывается тем, с чем само же сравнивается."""
    baselines = {}
    for key, entry in history.get("channels", {}).items():
        hist = [r for r in entry["history"] if r["date"] not in exclude_dates]
        n = len(hist)
        b7, b28 = hist[-7:], hist[-28:]
        last14 = hist[-14:]
        prior14 = hist[-28:-14] if len(hist) >= 28 else []

        metrics_out = {}
        for m in HISTORY_METRICS:
            mean7 = _mean([r[m] for r in b7]) if b7 else None
            mean28 = _mean([r[m] for r in b28]) if b28 else None
            std28 = _sample_std([r[m] for r in b28]) if len(b28) >= BASELINE_MIN_N else None
            trend = None
            if len(hist) >= 28 and prior14:
                mean_prior14 = _mean([r[m] for r in prior14])
                if mean_prior14:
                    trend = (_mean([r[m] for r in last14]) - mean_prior14) / mean_prior14 * 100
            metrics_out[m] = {"mean_7d": mean7, "mean_28d": mean28, "std_28d": std28, "trend_14v14_pct": trend}

        cost_28d_sum = sum(r["cost"] for r in b28)
        leads_28d_sum = sum(r["leads"] for r in b28)
        cost_7d_sum = sum(r["cost"] for r in b7)
        leads_7d_sum = sum(r["leads"] for r in b7)
        confidence = "low" if n < 14 else ("medium" if n < 42 else "high")
        baselines[key] = {
            "history_days_available": n,
            "baseline_confidence": confidence,
            "metrics": metrics_out,
            "cpl_pooled": {
                "value_28d": (cost_28d_sum / leads_28d_sum) if leads_28d_sum else None,
                "value_7d": (cost_7d_sum / leads_7d_sum) if leads_7d_sum else None,
                "leads_n_28d": leads_28d_sum,
                "low_volume_flag": (leads_28d_sum / 28.0) < 3,
            },
        }
    return baselines


def _z_score(current, mean, std):
    if current is None or mean is None or std is None:
        return None
    eff_std = max(std, STD_RELATIVE_FLOOR * abs(mean), STD_ABS_EPSILON)
    return (current - mean) / eff_std if eff_std else None


def select_top_channels_for_prompt(baselines, current_channels, n_full=5, n_total=20):
    """Отдаём Claude компактную, ограниченную по размеру сводку: полная
    детализация только для топ-5 каналов по вчерашнему расходу, урезанная
    для мест 6-20, остальное — одной агрегированной строкой. Так промпт не
    растёт линейно с числом мелких каналов/под-источников."""
    ranked = sorted(current_channels, key=lambda c: c["cost"], reverse=True)
    out = []
    other = {"channel": "__other_channels__", "channels_aggregated_count": 0, "visits": 0, "leads": 0, "cost": 0}
    for i, cur in enumerate(ranked):
        b = baselines.get(cur["key"])
        if not b:
            b = {"history_days_available": 0, "baseline_confidence": "low", "metrics": {}, "cpl_pooled": {}}
        if i < n_full:
            metrics_with_z = {}
            for m in HISTORY_METRICS:
                bm = b["metrics"].get(m, {})
                metrics_with_z[m] = {**bm, "yesterday": cur.get(m), "z_28d": _z_score(cur.get(m), bm.get("mean_28d"), bm.get("std_28d"))}
            out.append({
                "channel": cur["channel"],
                "history_days_available": b["history_days_available"],
                "baseline_confidence": b["baseline_confidence"],
                "metrics": metrics_with_z,
                "cpl_pooled": {**b["cpl_pooled"], "yesterday_cpl": cur.get("cpl")},
            })
        elif i < n_total:
            out.append({
                "channel": cur["channel"],
                "baseline_confidence": b["baseline_confidence"],
                "yesterday_leads": cur.get("leads"),
                "yesterday_cost": cur.get("cost"),
                "cpl_28d": b["cpl_pooled"].get("value_28d"),
            })
        else:
            other["channels_aggregated_count"] += 1
            other["visits"] += cur.get("visits", 0)
            other["leads"] += cur.get("leads", 0)
            other["cost"] += cur.get("cost", 0)
    if other["channels_aggregated_count"] > 0:
        out.append(other)
    return out


# ---------- 2. База знаний ----------
# knowledge_base.md уже давно существует и заявляет о себе, что его читает
# скрипт — но раньше в коде физически не было ни одной операции чтения этого
# файла. Плюс: подгрузка папки knowledge/ (PDF/PNG/JPG уходят к Claude как
# document/image-блоки; DOCX/PPTX/XLSX — текстом), как и было изначально
# задумано (пакеты python-pptx/python-docx/openpyxl не просто так стоят в
# workflow).
def load_knowledge():
    text_parts = []
    content_blocks = []

    if os.path.exists(KNOWLEDGE_BASE_PATH):
        with open(KNOWLEDGE_BASE_PATH, encoding="utf-8") as f:
            md = f.read().strip()
        if md:
            text_parts.append(f"### {KNOWLEDGE_BASE_PATH}\n{md}")

    if os.path.isdir(KNOWLEDGE_DIR):
        for name in sorted(os.listdir(KNOWLEDGE_DIR)):
            path = os.path.join(KNOWLEDGE_DIR, name)
            if not os.path.isfile(path):
                continue
            ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""
            try:
                if ext == "docx":
                    from docx import Document
                    doc = Document(path)
                    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                    if text:
                        text_parts.append(f"### {name}\n{text}")
                elif ext == "pptx":
                    from pptx import Presentation
                    prs = Presentation(path)
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        parts = [
                            shape.text_frame.text
                            for shape in slide.shapes
                            if shape.has_text_frame and shape.text_frame.text.strip()
                        ]
                        if parts:
                            slides_text.append(f"Слайд {i}: " + " / ".join(parts))
                    if slides_text:
                        text_parts.append(f"### {name}\n" + "\n".join(slides_text))
                elif ext == "xlsx":
                    from openpyxl import load_workbook
                    wb = load_workbook(path, data_only=True)
                    rows_text = []
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            vals = [str(v) for v in row if v is not None]
                            if vals:
                                rows_text.append(" | ".join(vals))
                    if rows_text:
                        text_parts.append(f"### {name}\n" + "\n".join(rows_text[:300]))
                elif ext == "pdf":
                    with open(path, "rb") as fbin:
                        content_blocks.append({
                            "type": "document",
                            "source": {
                                "type": "base64",
                                "media_type": "application/pdf",
                                "data": base64.b64encode(fbin.read()).decode(),
                            },
                            "title": name,
                        })
                elif ext in ("png", "jpg", "jpeg"):
                    media_type = "image/png" if ext == "png" else "image/jpeg"
                    with open(path, "rb") as fbin:
                        content_blocks.append({
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": base64.b64encode(fbin.read()).decode(),
                            },
                        })
            except Exception as e:
                print(f"[База знаний] Не удалось прочитать {name}: {e}")

    return "\n\n".join(text_parts), content_blocks


# ---------- 3. Claude: структурированный анализ ----------
# Раньше ответ Claude был свободным текстом с max_tokens=1500 — подтверждено
# (stop_reason="max_tokens"), что он обрывался ДО блоков "Что делать" и
# "Вопросы агентству". Теперь анализ идёт через принудительный tool-call:
# это гарантирует, что нужные поля либо есть все, либо вызов явно падает
# (и это видно в логе), вместо тихого обрезания текста на полуслове.
ANALYSIS_TOOL = {
    "name": "submit_analysis",
    "description": "Структурированный маркетинговый разбор периода для дашборда KSM.",
    "input_schema": {
        "type": "object",
        "properties": {
            "what_happened": {
                "type": "string",
                "description": "Что произошло — сравнение периода с предыдущим, коротко и по цифрам.",
            },
            "why": {
                "type": "string",
                "description": "Вероятные причины роста/падения CPL и других ключевых метрик.",
            },
            "actions": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Конкретные рекомендации: какие кампании/каналы отключить или усилить, какие ключевые слова/объявления/посадочные проверить.",
            },
            "agency_questions": {
                "type": "array",
                "items": {"type": "string"},
                "description": "2-4 вопроса рекламному агентству по сути отклонений.",
            },
        },
        "required": ["what_happened", "why", "actions", "agency_questions"],
    },
}


def _coerce_str_list(value, tag_name):
    """Claude почти всегда возвращает нормальный JSON-массив строк для полей
    actions/agency_questions — но на практике (подтверждено на реальном
    прогоне с большим промптом) иногда вместо этого приходит одна строка с
    псевдо-XML-тегами <action>...</action>/<question>...</question> внутри.
    Без этой нормализации такое значение долетало бы до сайта как есть и
    ломало рендер (список ожидает массив, а не строку). Пытаемся вытащить
    пункты из тегов; если не получилось — хотя бы не теряем текст целиком."""
    if isinstance(value, list):
        return [str(v) for v in value]
    if isinstance(value, str):
        import re
        matches = re.findall(rf"<{tag_name}>(.*?)</{tag_name}>", value, re.S)
        if matches:
            return [m.strip() for m in matches if m.strip()]
        cleaned = value.strip()
        return [cleaned] if cleaned else []
    return []


def call_claude_structured(system_prompt, user_text, content_blocks=None):
    content = [{"type": "text", "text": user_text}]
    if content_blocks:
        content = content + content_blocks
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers={
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01",
            "content-type": "application/json",
        },
        json={
            "model": CLAUDE_MODEL,
            "max_tokens": CLAUDE_MAX_TOKENS,
            "system": system_prompt,
            "tools": [ANALYSIS_TOOL],
            "tool_choice": {"type": "tool", "name": "submit_analysis"},
            "messages": [{"role": "user", "content": content}],
        },
        timeout=90,
    )
    resp.raise_for_status()
    body = resp.json()
    stop_reason = body.get("stop_reason")
    if stop_reason == "max_tokens":
        print(f"[Claude] ВНИМАНИЕ: ответ оборвался по max_tokens ({CLAUDE_MAX_TOKENS}) — возможна неполная структура.")
    for block in body.get("content", []):
        if block.get("type") == "tool_use" and block.get("name") == "submit_analysis":
            result = block["input"]
            for field, tag in (("actions", "action"), ("agency_questions", "question")):
                original = result.get(field)
                if not isinstance(original, list):
                    print(f"[Claude] Поле «{field}» пришло не списком ({type(original).__name__}), нормализую.")
                    result[field] = _coerce_str_list(original, tag)
            return result
    raise RuntimeError(f"Claude не вернул структурированный анализ (tool_use не найден): {body}")


DAILY_SYSTEM_PROMPT_TEMPLATE = """Ты — ИИ-агент «Маркетинг-аналитик KSM» для отдела маркетинга компании,
производящей строительные материалы (Казахстан, валюта — тенге ₸).
Тебе присылают уже посчитанные кодом (не тобой) реальные показатели рекламы
из Roistat за вчера и позавчера: расход, лиды, CPL, визиты, показы, клики,
CTR, CPC и разбивку по каналам. Используй именно эти цифры, не пересчитывай
и не придумывай новые.

База знаний компании (используй при выводах и рекомендациях, если это
уместно):
{knowledge}

Заполни структурированный анализ через инструмент submit_analysis на
русском, кратко и по делу, не придумывая цифр, которых нет во входных
данных."""

WEEKLY_SYSTEM_PROMPT_TEMPLATE = """Ты — ИИ-агент «Маркетинг-аналитик KSM» для отдела маркетинга компании,
производящей строительные материалы (Казахстан, валюта — тенге ₸).
Тебе присылают уже посчитанные кодом реальные показатели рекламы из Roistat
за ЗАВЕРШЁННУЮ РАБОЧУЮ НЕДЕЛЮ (со среды по вторник) и за такую же
предыдущую неделю: расход, лиды, CPL, визиты, показы, клики, CTR, CPC и
разбивку по каналам. Используй именно эти цифры, не пересчитывай и не
придумывай новые.

База знаний компании (используй при выводах и рекомендациях, если это
уместно):
{knowledge}

Заполни структурированный анализ через инструмент submit_analysis на
русском: опиши недельную динамику, вероятные причины, бюджетные
рекомендации на неделю вперёд и вопросы агентству по итогам недели."""


# ---------- 4. Telegram ----------
# Самостоятельная подписка: сотрудник пишет боту, бот сам спрашивает кодовое
# слово и регистрирует новый chat_id после того, как слово названо — без
# ручного вмешательства. Так как скрипт запускается не постоянно, а раз в
# день по расписанию, разговор растягивается на реальных запусках: "здесь
# спросили код" и "тут его назвали" могут оказаться в разных запусках, но
# итог тот же — подписка без правки секретов и кода.
RECIPIENTS_PATH = os.path.join(DATA_DIR, "telegram_recipients.json")


def load_recipients():
    if not os.path.exists(RECIPIENTS_PATH):
        return {"chat_ids": [], "pending": [], "last_update_id": 0}
    try:
        with open(RECIPIENTS_PATH, encoding="utf-8") as f:
            data = json.load(f)
        data.setdefault("chat_ids", [])
        data.setdefault("pending", [])
        data.setdefault("last_update_id", 0)
        return data
    except (json.JSONDecodeError, OSError) as e:
        print(f"[Telegram] Не удалось прочитать {RECIPIENTS_PATH}, начинаю с чистого листа: {e}", file=sys.stderr)
        return {"chat_ids": [], "pending": [], "last_update_id": 0}


def save_recipients(data):
    os.makedirs(DATA_DIR, exist_ok=True)
    with open(RECIPIENTS_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _send_single_telegram(chat_id, text):
    if not TELEGRAM_BOT_TOKEN:
        return
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage"
    try:
        resp = requests.post(url, json={"chat_id": chat_id, "text": text}, timeout=15)
        if resp.status_code >= 400:
            print(f"[Telegram] Ошибка отправки на chat_id {chat_id}, статус {resp.status_code}: {resp.text[:300]}", file=sys.stderr)
    except requests.exceptions.RequestException as e:
        print(f"[Telegram] Сетевая ошибка при отправке на chat_id {chat_id}: {e}", file=sys.stderr)


def sync_telegram_recipients():
    """Проверяет новые сообщения боту (getUpdates) и ведёт саморегистрацию
    по кодовому слову TELEGRAM_JOIN_CODE. Возвращает актуальный список
    chat_id, которые уже подписаны (из data/telegram_recipients.json)."""
    if not TELEGRAM_BOT_TOKEN:
        return []
    data = load_recipients()
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getUpdates"
    try:
        resp = requests.get(url, params={"offset": data["last_update_id"] + 1}, timeout=15)
        resp.raise_for_status()
        updates = resp.json().get("result", [])
    except requests.exceptions.RequestException as e:
        print(f"[Telegram] Не удалось получить обновления для авторегистрации: {e}", file=sys.stderr)
        return data["chat_ids"]

    latest_text_by_chat = {}
    for upd in updates:
        data["last_update_id"] = max(data["last_update_id"], upd.get("update_id", 0))
        msg = upd.get("message")
        if not msg or (msg.get("chat") or {}).get("type") != "private":
            continue
        chat_id = str(msg["chat"]["id"])
        latest_text_by_chat[chat_id] = (msg.get("text") or "").strip().lower()

    for chat_id, text in latest_text_by_chat.items():
        if chat_id in data["chat_ids"]:
            continue
        code_matched = bool(TELEGRAM_JOIN_CODE) and TELEGRAM_JOIN_CODE in text
        if code_matched:
            data["chat_ids"].append(chat_id)
            if chat_id in data["pending"]:
                data["pending"].remove(chat_id)
            _send_single_telegram(chat_id, "Готово! Вы подписаны на ежедневные и еженедельные отчёты KSM-маркетинга.")
            print(f"[Telegram] Новый подписчик через кодовое слово: {chat_id}")
        elif chat_id in data["pending"]:
            _send_single_telegram(chat_id, "Кодовое слово не распознано. Напишите его ещё раз, чтобы подписаться на отчёты KSM-маркетинга.")
        else:
            data["pending"].append(chat_id)
            _send_single_telegram(chat_id, "Здравствуйте! Чтобы подписаться на отчёты KSM-маркетинга, напишите кодовое слово.")

    save_recipients(data)
    return data["chat_ids"]


def send_telegram_alert(text):
    if not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_IDS:
        print("Telegram не настроен (нет токена/chat_id) — пропускаю отправку.")
        return
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage"
    # Рассылка каждому получателю по отдельности — если один chat_id окажется
    # неверным/отозванным, это не должно останавливать доставку остальным.
    for chat_id in TELEGRAM_CHAT_IDS:
        resp = requests.post(url, json={"chat_id": chat_id, "text": text}, timeout=15)
        if resp.status_code >= 400:
            print(f"[Telegram] Ошибка отправки на chat_id {chat_id}, статус {resp.status_code}: {resp.text[:300]}")
        else:
            print(f"[Telegram] Сообщение отправлено (chat_id {chat_id}).")


def send_daily_telegram(snapshot):
    """Отправляется КАЖДЫЙ день, а не только когда сработала тревога —
    руководителю нужен ежедневный отчёт, а не только сигналы проблем."""
    y = snapshot["totals"]["yesterday"]
    lines = [
        f"KSM-маркетинг, отчёт за {snapshot['period']['yesterday']}:",
        f"Лиды: {y['leads']}, расход: {y['cost']:.0f} ₸"
        + (f", CPL: {y['cpl']:.0f} ₸." if y.get("cpl") else "."),
    ]
    if snapshot.get("best_channel"):
        b = snapshot["best_channel"]
        lines.append(f"Лучший канал: {b['channel']} (CPL {b['cpl']:.0f} ₸).")
    if snapshot.get("worst_channel") and snapshot["worst_channel"] is not snapshot.get("best_channel"):
        w = snapshot["worst_channel"]
        lines.append(f"Худший канал: {w['channel']} (CPL {w['cpl']:.0f} ₸).")
    if snapshot["alerts"]:
        lines.append("")
        lines.append("Внимание:")
        lines += [f"- {a}" for a in snapshot["alerts"]]
    lines.append("")
    lines.append(snapshot["analysis"]["what_happened"][:600])
    send_telegram_alert("\n".join(lines))


def send_weekly_telegram(snapshot):
    tw = snapshot["totals"]["this_week"]
    pw = snapshot["totals"]["prev_week"]
    p = snapshot["period"]
    lines = [
        f"KSM-маркетинг, недельный отчёт ({p['week_start']} — {p['week_end']}):",
        f"Лиды: {tw['leads']} (было {pw['leads']}), "
        f"расход: {tw['cost']:.0f} ₸ (было {pw['cost']:.0f} ₸)"
        + (f", CPL: {tw['cpl']:.0f} ₸." if tw.get("cpl") else "."),
    ]
    if snapshot["alerts"]:
        lines.append("")
        lines.append("Отклонения за неделю:")
        lines += [f"- {a}" for a in snapshot["alerts"]]
    lines.append("")
    lines.append(snapshot["analysis"]["what_happened"][:600])
    send_telegram_alert("\n".join(lines))


# ---------- 5. Проверка порогов (без ИИ, для мгновенного сигнала) ----------
# Раньше здесь было только 2 сигнала (падение лидов, рост CPL). ТЗ требует 8
# типов уведомлений — остальные 6 добавлены ниже там, где это физически
# возможно посчитать по данным, которые реально отдаёт Roistat.
def check_thresholds(curr, prev, curr_channels, cfg=None):
    cfg = cfg or CONFIG
    alerts = []

    if prev.get("leads", 0) > 0:
        change = (curr["leads"] - prev["leads"]) / prev["leads"] * 100
        if change <= -cfg["leadsDropAlertPct"]:
            alerts.append(f"Лиды упали на {abs(change):.1f}% ({curr['leads']} против {prev['leads']}).")

    if curr.get("cpl") and prev.get("cpl"):
        change = (curr["cpl"] - prev["cpl"]) / prev["cpl"] * 100
        if change >= cfg["cplAlertPct"]:
            alerts.append(f"CPL вырос на {change:.1f}% (с {prev['cpl']:.0f} до {curr['cpl']:.0f} ₸).")

    if prev.get("cost", 0) > 0:
        cost_change = (curr["cost"] - prev["cost"]) / prev["cost"] * 100
        if cost_change >= cfg["cplAlertPct"] and curr["leads"] <= prev["leads"]:
            alerts.append(
                f"Расход вырос на {cost_change:.1f}%, а заявок больше не стало "
                f"({curr['leads']} против {prev['leads']})."
            )

    if prev.get("visits", 0) > 0:
        change = (curr["visits"] - prev["visits"]) / prev["visits"] * 100
        if change <= -cfg["visitsDropAlertPct"]:
            alerts.append(f"Трафик резко снизился на {abs(change):.1f}% ({curr['visits']} визитов против {prev['visits']}).")

    if curr.get("ctr") is not None and prev.get("ctr"):
        change = (curr["ctr"] - prev["ctr"]) / prev["ctr"] * 100
        if change <= -cfg["ctrDropAlertPct"]:
            alerts.append(f"CTR упал на {abs(change):.1f}% (с {prev['ctr']:.2f}% до {curr['ctr']:.2f}%).")

    if curr.get("cpc") is not None and prev.get("cpc"):
        change = (curr["cpc"] - prev["cpc"]) / prev["cpc"] * 100
        if change >= cfg["cpcRiseAlertPct"]:
            alerts.append(f"CPC вырос на {change:.1f}% (с {prev['cpc']:.0f} до {curr['cpc']:.0f} ₸).")

    if curr.get("cr") is not None and prev.get("cr"):
        change = (curr["cr"] - prev["cr"]) / prev["cr"] * 100
        if change <= -cfg["conversionDropAlertPct"]:
            alerts.append(f"Конверсия визитов в лиды упала на {abs(change):.1f}% (с {prev['cr']:.1f}% до {curr['cr']:.1f}%).")

    for ch in curr_channels:
        if ch["leads"] == 0 and ch["cost"] >= cfg["noLeadsSpendAlertTenge"]:
            alerts.append(f"«{ch['channel']}» потратил {ch['cost']:.0f} ₸ без единого лида.")

    return alerts


def check_plan_progress(month_totals, days_elapsed, days_in_month, cfg=None):
    """Факт к KPI/плану — раздел 5 и пункт 7 раздела 8 ТЗ. План задаётся в
    config.json (monthlyLeadsTarget/monthlyBudgetTarget), редактируется на
    сайте во вкладке «Настройки»."""
    cfg = cfg or CONFIG
    alerts = []
    if days_in_month <= 0:
        return alerts
    expected_leads = cfg["monthlyLeadsTarget"] * days_elapsed / days_in_month
    if expected_leads > 0 and month_totals["leads"] < expected_leads * 0.8:
        alerts.append(
            f"Лиды с начала месяца отстают от плана: {month_totals['leads']} при ожидаемых "
            f"~{expected_leads:.0f} на {days_elapsed}-й день месяца (план на весь месяц — "
            f"{cfg['monthlyLeadsTarget']})."
        )
    return alerts


def _days_in_month(d):
    if d.month == 12:
        next_month_first = datetime.date(d.year + 1, 1, 1)
    else:
        next_month_first = datetime.date(d.year, d.month + 1, 1)
    return (next_month_first - timedelta(days=1)).day


# ---------- 6. Неделя KSM: среда → вторник ----------
def get_last_completed_week_bounds(today):
    """Возвращает (начало, конец) последней ПОЛНОСТЬЮ завершённой недели
    среда-вторник относительно today. Вызывается по средам — тогда конец
    недели это вчера (вторник), начало — среда 6 дней до этого."""
    # weekday(): Monday=0 ... Sunday=6. Вторник=1, среда=2.
    days_since_tuesday = (today.weekday() - 1) % 7
    if days_since_tuesday == 0:
        days_since_tuesday = 7
    week_end = today - timedelta(days=days_since_tuesday)
    week_start = week_end - timedelta(days=6)
    return week_start, week_end


def is_weekly_report_day(today):
    return today.weekday() == 2  # среда


def run_weekly_report(today, knowledge_text, knowledge_blocks, history):
    week_start, week_end = get_last_completed_week_bounds(today)
    prev_week_end = week_start - timedelta(days=1)
    prev_week_start = prev_week_end - timedelta(days=6)

    raw_week = fetch_roistat_data(str(week_start), str(week_end))
    raw_prev_week = fetch_roistat_data(str(prev_week_start), str(prev_week_end))

    totals_week = extract_totals(raw_week)
    totals_prev_week = extract_totals(raw_prev_week)
    channels_week = extract_channels(raw_week)
    best_channel, worst_channel, no_leads_channels = rank_channels(channels_week)

    alerts = check_thresholds(totals_week, totals_prev_week, channels_week)

    # Те же накопленные базовые линии, что и в дневном отчёте (90-дневное
    # окно из data/channel_history.json), исключая обе недели этого
    # сравнения — недельный отчёт не заменяет отдельного истории-файла,
    # использует тот же самый.
    week_dates = {str(week_start + timedelta(days=i)) for i in range(7)}
    prev_week_dates = {str(prev_week_start + timedelta(days=i)) for i in range(7)}
    weekly_baselines = compute_baselines(history, exclude_dates=week_dates | prev_week_dates)
    channel_baselines = select_top_channels_for_prompt(weekly_baselines, channels_week)

    user_text = (
        f"Неделя {week_start}..{week_end} против недели {prev_week_start}..{prev_week_end} "
        "(посчитано кодом, не пересчитывать):\n"
        + json.dumps(
            {
                "эта_неделя": totals_week,
                "прошлая_неделя": totals_prev_week,
                "каналы_эта_неделя": channels_week,
                "накопленная_статистика_по_каналам": channel_baselines,
            },
            ensure_ascii=False,
        )
        + "\n\nПоле «накопленная_статистика_по_каналам» — среднее/std/тренд по каждому каналу "
        "за последние 7-28 дней (не считая обе недели сравнения) и z_28d — отклонение "
        "недельного значения от обычной нормы канала в стандартных отклонениях."
    )
    system_prompt = WEEKLY_SYSTEM_PROMPT_TEMPLATE.format(knowledge=knowledge_text or "(база знаний пуста)")
    analysis = call_claude_structured(system_prompt, user_text, knowledge_blocks)

    weekly_snapshot = {
        "updated_at": datetime.datetime.now(timezone.utc).isoformat(),
        "period": {
            "week_start": str(week_start),
            "week_end": str(week_end),
            "prev_week_start": str(prev_week_start),
            "prev_week_end": str(prev_week_end),
        },
        "totals": {"this_week": totals_week, "prev_week": totals_prev_week},
        "channels": {"this_week": channels_week},
        "channel_baselines": channel_baselines,
        "best_channel": best_channel,
        "worst_channel": worst_channel,
        "no_leads_spend_channels": no_leads_channels,
        "analysis": analysis,
        "alerts": alerts,
    }
    with open(WEEKLY_SNAPSHOT_PATH, "w", encoding="utf-8") as f:
        json.dump(weekly_snapshot, f, ensure_ascii=False, indent=2)

    send_weekly_telegram(weekly_snapshot)


# ---------- main: одна ежедневная итерация ----------
def main():
    global TELEGRAM_CHAT_IDS
    today = almaty_today()
    yesterday_date = today - timedelta(days=1)
    day_before_date = today - timedelta(days=2)

    # Сначала синхронизация подписчиков (кодовое слово в боте), чтобы те, кто
    # уже подписался, сразу получили и сегодняшний отчёт, а не только со
    # следующего запуска.
    dynamic_recipients = sync_telegram_recipients()
    TELEGRAM_CHAT_IDS = sorted(set(TELEGRAM_CHAT_IDS) | set(dynamic_recipients))

    knowledge_text, knowledge_blocks = load_knowledge()

    raw_yesterday = fetch_roistat_data(str(yesterday_date), str(yesterday_date))
    raw_day_before = fetch_roistat_data(str(day_before_date), str(day_before_date))

    totals_yesterday = extract_totals(raw_yesterday)
    totals_day_before = extract_totals(raw_day_before)
    channels_yesterday = extract_channels(raw_yesterday)
    channels_day_before = extract_channels(raw_day_before)
    best_channel, worst_channel, no_leads_channels = rank_channels(channels_yesterday)

    alerts = check_thresholds(totals_yesterday, totals_day_before, channels_yesterday)

    # Факт к плану/KPI (месяц-к-дате) — раздел 5 и пункт 7 раздела 8 ТЗ.
    month_start = yesterday_date.replace(day=1)
    raw_month = fetch_roistat_data(str(month_start), str(yesterday_date))
    month_totals = extract_totals(raw_month)
    days_elapsed = (yesterday_date - month_start).days + 1
    days_in_month = _days_in_month(yesterday_date)
    alerts += check_plan_progress(month_totals, days_elapsed, days_in_month)

    # Накопительная статистика по каналам: считаем базовые линии ДО того, как
    # добавим сегодняшнюю запись в историю (иначе день сравнивался бы сам с
    # собой). Исключаем и вчера, и позавчера — оба уже участвуют в дневном
    # сравнении напрямую.
    history = load_history()
    baselines = compute_baselines(history, exclude_dates={str(yesterday_date), str(day_before_date)})
    channel_baselines = select_top_channels_for_prompt(baselines, channels_yesterday)

    daily_user_text = (
        "Данные за вчера и позавчера (посчитано кодом, не пересчитывать):\n"
        + json.dumps(
            {
                "вчера": totals_yesterday,
                "позавчера": totals_day_before,
                "каналы_вчера": channels_yesterday,
                "накопленная_статистика_по_каналам": channel_baselines,
            },
            ensure_ascii=False,
        )
        + "\n\nПоле «накопленная_статистика_по_каналам» — это среднее/std/тренд по каждому "
        "каналу за последние 7-28 дней (сколько накопилось) и z_28d — насколько вчерашнее "
        "значение отклоняется от нормы канала в единицах стандартного отклонения. "
        "baseline_confidence показывает, сколько дней истории накоплено (low/medium/high) — "
        "при low не делай категоричных выводов, данных ещё мало."
    )
    daily_system_prompt = DAILY_SYSTEM_PROMPT_TEMPLATE.format(knowledge=knowledge_text or "(база знаний пуста)")
    analysis = call_claude_structured(daily_system_prompt, daily_user_text, knowledge_blocks)

    daily_snapshot = {
        "updated_at": datetime.datetime.now(timezone.utc).isoformat(),
        "period": {"today": str(today), "yesterday": str(yesterday_date), "day_before": str(day_before_date)},
        "totals": {"yesterday": totals_yesterday, "day_before": totals_day_before},
        "channels": {"yesterday": channels_yesterday, "day_before": channels_day_before},
        "channel_baselines": channel_baselines,
        "best_channel": best_channel,
        "worst_channel": worst_channel,
        "no_leads_spend_channels": no_leads_channels,
        "analysis": analysis,
        "alerts": alerts,
        "month_to_date": {
            "totals": month_totals,
            "days_elapsed": days_elapsed,
            "days_in_month": days_in_month,
            "leads_target": CONFIG["monthlyLeadsTarget"],
            "budget_target": CONFIG["monthlyBudgetTarget"],
        },
    }
    with open(DAILY_SNAPSHOT_PATH, "w", encoding="utf-8") as f:
        json.dump(daily_snapshot, f, ensure_ascii=False, indent=2)

    send_daily_telegram(daily_snapshot)

    if is_weekly_report_day(today):
        run_weekly_report(today, knowledge_text, knowledge_blocks, history)

    # Историю пополняем последним шагом — после того, как baseline для
    # сегодняшнего сравнения уже посчитан без сегодняшней записи. Перезабираем
    # не только "вчера", а последние REFRESH_WINDOW_DAYS дней — Roistat может
    # задним числом пересмотреть уже отданные цифры (подтверждено на практике).
    refresh_recent_history(history, today)
    save_history(history)

    print("Готово:", datetime.datetime.now().isoformat())


if __name__ == "__main__":
    main()
